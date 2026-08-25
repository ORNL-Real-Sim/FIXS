#!/usr/bin/env python
"""Regenerate CarlaDynoCoupling.pptx.

Companion deck to AxleDynoCarMakerCoupling.pptx. Same visual system -- the drawing
primitives and colour tokens are imported from make_slides_pptx rather than copied,
so the two decks cannot drift apart.

    pip install python-pptx
    python doc/xil/make_carla_slides_pptx.py

Content follows CarlaDynoCoupling.md: what CARLA will not tell you, the two designs
that route around it, the ceiling they share, and what is measured versus derived.
"""
from __future__ import annotations

import pathlib
import sys

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent))

from make_slides_pptx import (  # noqa: E402
    DASH, F_BODY, F_DISPLAY, F_MONO,
    HW, HW_SOFT, INK, INK_2, INK_3, MODEL, MODEL_SOFT, PANEL, PAPER, RULE,
    SH, SW, WARN, WARN_SOFT,
    arrow, box, caption, inch, labelled_box, line, mono_block, plain_line, textbox,
)

OUT = pathlib.Path(__file__).resolve().parent / "CarlaDynoCoupling.pptx"
NSLIDES = 6


def new_slide(prs, num, eyebrow, title, lede=None):
    """Same header as the CarMaker deck, with this deck's slide count."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER

    plain_line(s, 0.62, 0.52, 12.72, 0.52, color=RULE, width=1.0)

    tf = textbox(s, 0.62, 0.24, 8.0, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "%02d / %02d" % (num, NSLIDES)
    r.font.name, r.font.size, r.font.bold = F_MONO, Pt(9), True
    r.font.color.rgb = HW
    r = p.add_run(); r.text = "    " + eyebrow.upper()
    r.font.name, r.font.size = F_MONO, Pt(9)
    r.font.color.rgb = INK_3

    tf = textbox(s, 0.62, 0.70, 12.1, 0.62)
    line(tf, title, font=F_DISPLAY, size=30, color=INK, bold=True, first=True)

    if lede:
        tf = textbox(s, 0.62, 1.36, 11.4, 0.5)
        line(tf, lede, font=F_BODY, size=12.5, color=INK_2, first=True)
    return s


# ---------------------------------------------------------------- 01
def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = PAPER

    tf = textbox(s, 0.9, 1.15, 8.0, 0.3)
    line(tf, "FIXS  ·  XIL BENCH  ·  DESIGN STUDY", font=F_MONO, size=10,
         color=HW, bold=True, first=True)

    tf = textbox(s, 0.9, 1.62, 11.2, 1.5)
    line(tf, "Axle Dyno ↔ CARLA", font=F_DISPLAY, size=52, color=INK, bold=True, first=True)

    tf = textbox(s, 0.9, 3.05, 9.8, 1.0)
    line(tf, "The dyno loop needs one signal from the vehicle model. CarMaker puts it "
             "on the wire. CARLA will not, and no configuration changes that — so the "
             "two designs here are two ways of routing around it.",
         font=F_BODY, size=14, color=INK_2, first=True)

    mono_block(s, 0.9, 4.25, 8.2, 0.72, [
        ("I * dw/dt  =  Taxl  +  Trq_T2W", INK, True),
        ("                          ^^^^^^^  CARLA will not give you this", HW, False),
    ])

    tf = textbox(s, 0.9, 5.30, 11.0, 0.56)
    for t in ["Design A   impose the speed      dyno owns longitudinal, CARLA is told the answer",
              "Design B   inject the torque     CARLA's tire and body run, speed recovered from motion"]:
        line(tf, t, font=F_MONO, size=11, color=INK_2, first=t.startswith("Design A"))
    tf = textbox(s, 0.9, 5.92, 11.0, 0.3)
    line(tf, "Both are lumped single-axis. Per-wheel needs a fork.",
         font=F_MONO, size=11, color=WARN, bold=True, first=True)
    return s


# ---------------------------------------------------------------- 02
def slide_02(prs):
    s = new_slide(prs, 2, "The wall — on 0.9.15",
                  "What CARLA 0.9.15 will and will not tell you",
                  "The data exists and is computed every tick. On 0.9.15 there is no way "
                  "to ask for it. On 0.9.16 there is — see slide 06.")

    labelled_box(s, 0.62, 2.00, 5.9, 2.05, "READABLE",
                 ["get_velocity          get_transform",
                  "get_angular_velocity  get_control",
                  "get_acceleration      get_physics_control",
                  "",
                  "get_wheel_steer_angle    <- the only",
                  "                            per-wheel getter"],
                 fill=MODEL_SOFT, edge=MODEL, head_color=MODEL, width=1.5)

    labelled_box(s, 6.82, 2.00, 5.9, 2.05, "NOT READABLE",
                 ["wheel speed        omega_i",
                  "longitudinal slip  kappa_i",
                  "tire force         Fx_i",
                  "normal load        Fz_i",
                  "wheel torque       Trq_T2W",
                  ""],
                 fill=WARN_SOFT, edge=WARN, head_color=WARN, width=2.0)

    tf = textbox(s, 0.62, 4.28, 12.1, 0.3)
    line(tf, "but all of it is already computed and stored:", font=F_BODY,
         size=12, color=INK_2, first=True)

    mono_block(s, 0.62, 4.62, 12.1, 1.16, [
        ("UVehicleWheel  (UE4)   DebugLongSlip   DebugLatSlip    DebugTireLoad", INK, False),
        ("                       DebugWheelTorque  DebugLongForce  DebugLatForce", INK, False),
        ("CARLA already holds the pointer: CarlaWheeledVehicle.cpp:309", INK_3, False),
        ("show_debug_telemetry renders them to the HUD.  There is no RPC path.", HW, True),
    ])

    caption(s, "Consequence on 0.9.15: Trq_T2W and omega_i are both unavailable, and the "
               "two designs that follow are ways of living with that.", y=5.98, color=WARN)
    return s


# ---------------------------------------------------------------- 03
def slide_03(prs):
    s = new_slide(prs, 3, "Design A", "Impose the speed",
                  "The dyno owns all longitudinal dynamics, including road load from the "
                  "vehicle's own measured coefficients. CARLA is told the answer.")

    labelled_box(s, 0.62, 2.02, 3.5, 1.55, "DYNO + ROAD LOAD",
                 ["Taxl from the real vehicle", "A, B, C from real coastdown",
                  "-> v(k)"],
                 fill=HW_SOFT, edge=HW, head_color=HW, width=2.0)

    labelled_box(s, 4.92, 2.02, 3.5, 1.55, "CARLA",
                 ["enable_constant_velocity", "( v, 0, 0 )  body frame", "",
                  "pose · lateral · traffic"],
                 fill=MODEL_SOFT, edge=MODEL, head_color=MODEL, width=1.5)

    labelled_box(s, 9.22, 2.02, 3.5, 1.55, "HUB DYNOS",
                 ["w_cmd_i = v / r", "evenly to all units", "",
                  "longitudinal only"],
                 fill=None, edge=INK, head_color=INK, width=2.0)

    arrow(s, 4.12, 2.72, 4.86, 2.72, color=HW, width=2.5)
    arrow(s, 8.42, 2.72, 9.16, 2.72, color=MODEL, width=2.5)

    tf = textbox(s, 0.62, 3.80, 6.0, 0.28)
    line(tf, "MEASURED", font=F_DISPLAY, size=13, color=MODEL, bold=True, first=True)
    tf = textbox(s, 0.62, 4.10, 6.1, 1.5)
    for t in ["rejects any external force completely; the residual is exactly",
              "one integration step of it and does not accumulate.",
              "On a 20-degree slope the error was 0.069 m/s -- 0.46%.",
              "",
              "the wheels keep up:  tau = J_w*v/(r^2*Cs) = 3.4e-4 * v",
              "  10 m/s -> 3.4 ms    20 -> 6.8 ms    30 -> 10.2 ms",
              "  against a 20-50 ms tick, 3 to 15 time constants."]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2, first=t.startswith("rejects"))

    tf = textbox(s, 6.90, 3.80, 6.0, 0.28)
    line(tf, "WHAT IT COSTS", font=F_DISPLAY, size=13, color=WARN, bold=True, first=True)
    tf = textbox(s, 6.90, 4.10, 5.9, 1.5)
    for t in ["writing along the nose imposes ZERO SIDESLIP",
              "  measured 0.45 deg, against 21 deg for a free car",
              "  in the same corner.",
              "",
              "the car moves exactly where it points. Yaw is still",
              "dynamic, so the path is right -- what you lose is",
              "drift, oversteer, anything at the friction limit."]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2,
             first=t.startswith("writing"))

    caption(s, "Use enable_constant_velocity, not set_target_velocity: the former is latched "
               "and in body frame (held 15.00 m/s indefinitely,", y=5.86)
    caption(s, "within 1.3 deg of the nose through a 130 deg turn). The latter is a one-shot "
               "world-frame write. Both need physics ENABLED.", y=6.10)
    return s


# ---------------------------------------------------------------- 04
def slide_04(prs):
    s = new_slide(prs, 4, "Design B", "Inject the torque",
                  "CARLA's own tire, wheel and body physics run. The bench feeds torque "
                  "through the pedals; the wheel speed is recovered from the body's motion.")

    tf = textbox(s, 0.62, 1.96, 12.1, 0.28)
    line(tf, "FIRST: the drivetrain has to be flattened, or throttle is not a torque",
         font=F_DISPLAY, size=14, color=WARN, bold=True, first=True)

    b = box(s, 0.62, 2.30, 5.9, 1.20, fill=WARN_SOFT, edge=WARN, width=2.0)
    tf = b.text_frame
    line(tf, "STOCK", font=F_MONO, size=10.5, color=WARN, bold=True, first=True)
    for t in ["moi = 1.0 through final_ratio = 9",
              "  = 81 kg m^2 at the wheel = 592 kg",
              "plus clutch slip",
              "MEASURED: 47% spread across throttle"]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2, space_before=2)

    b = box(s, 6.82, 2.30, 5.9, 1.20, fill=MODEL_SOFT, edge=MODEL, width=2.0)
    tf = b.text_frame
    line(tf, "FLATTENED  (setup-time, once)", font=F_MONO, size=10.5,
         color=MODEL, bold=True, first=True)
    for t in ["moi 0.01 · clutch 1e4 · final 1.0",
              "one gear 1.0 · autobox off",
              "flat torque curve · damping 0",
              "MEASURED: 2.7% spread, 1.1% of predicted"]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2, space_before=2)

    arrow(s, 6.52, 2.90, 6.76, 2.90, color=INK, width=2.5)

    mono_block(s, 0.62, 3.68, 12.1, 0.86, [
        ("Taxl >= 0 :  throttle = Taxl / T0", HW, True),
        ("Taxl <  0 :  brake    = |Taxl| / ( n_wheels * max_brake_torque )", HW, True),
    ])

    tf = textbox(s, 0.62, 4.72, 12.1, 0.28)
    line(tf, "THEN: recover the wheel speed from the body",
         font=F_DISPLAY, size=14, color=INK, bold=True, first=True)

    mono_block(s, 0.62, 5.06, 12.1, 1.06, [
        ("sum Fx  =  m * dv/dt  +  Fres(v)  +  m*g*sin(theta)", INK, True),
        ("Trq_T2W =  -r * sum Fx", INK, True),
        ("Fres is the ONLY chassis force besides gravity -- everything else is inside Fx", INK_3, False),
    ])

    caption(s, "Difference get_velocity for dv/dt. Do NOT use get_acceleration -- it read "
               "15.55 m/s^2 vertically on a stationary car.", y=6.30, color=WARN)
    return s


# ---------------------------------------------------------------- 05
def slide_05(prs):
    s = new_slide(prs, 5, "The ceiling", "Both designs are lumped, single-axis",
                  "Not by choice. The public API constrains it at both ends, so a "
                  "per-wheel hub dyno gains nothing from Design B over Design A.")

    b = box(s, 0.62, 2.10, 5.9, 1.30, fill=None, edge=WARN, width=2.0)
    tf = b.text_frame
    line(tf, "IN", font=F_MONO, size=11, color=WARN, bold=True, first=True)
    for t in ["throttle is ONE scalar. Torque reaches the",
              "wheels through CARLA's differential, which",
              "decides the split.",
              "Left and right cannot be commanded apart."]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2, space_before=2)

    b = box(s, 6.82, 2.10, 5.9, 1.30, fill=None, edge=WARN, width=2.0)
    tf = b.text_frame
    line(tf, "OUT", font=F_MONO, size=11, color=WARN, bold=True, first=True)
    for t in ["the vdot inverse yields sum Fx from body",
              "acceleration. Splitting it per wheel needs",
              "per-wheel Fz and slip.",
              "Neither is exposed."]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2, space_before=2)

    tf = textbox(s, 0.62, 3.62, 12.1, 0.30)
    line(tf, "THE ONLY PER-WHEEL PATH IS A FORK — and the engine class already exists",
         font=F_DISPLAY, size=14, color=MODEL, bold=True, first=True)

    b = box(s, 0.62, 4.00, 12.1, 1.32, fill=MODEL_SOFT, edge=MODEL, width=2.0)
    tf = b.text_frame
    line(tf, "USimpleWheeledVehicleMovementComponent   ->   PxVehicleNoDrive",
         font=F_MONO, size=10.5, color=MODEL, bold=True, first=True)
    line(tf, '"a simple PxNoDrive which will give us suspension but no engine forces '
             'which we leave to the user"',
         font=F_MONO, size=9.5, color=INK_3, space_before=3)
    for t in ["void SetDriveTorque(float DriveTorque, int32 WheelIndex);",
              "void SetBrakeTorque(float BrakeTorque, int32 WheelIndex);",
              "void SetSteerAngle (float SteerAngle,  int32 WheelIndex);"]:
        line(tf, t, font=F_MONO, size=10, color=INK, space_before=2)

    caption(s, "Per-wheel torque, no engine, suspension and tires intact. What is missing is "
               "a CARLA vehicle class that uses it -- CARLA already", y=5.52)
    caption(s, "has the pattern, since ACarlaWheeledVehicleNW sits alongside the 4W one -- "
               "plus the RPC bindings. A scope decision, not a technical one.", y=5.76)
    return s


# ---------------------------------------------------------------- 06
def slide_06(prs):
    s = new_slide(prs, 6, "0.9.16 changes the answer",
                  "get_telemetry_data makes both signals readable",
                  "We are on 0.9.15.2. Everything on slides 02 and 04 exists to work "
                  "around a signal that the next release exposes directly.")

    b = box(s, 0.62, 1.96, 12.1, 1.62, fill=MODEL_SOFT, edge=MODEL, width=2.0)
    tf = b.text_frame
    line(tf, "carla.WheelTelemetryData      -- per wheel, read-only, no HUD needed",
         font=F_MONO, size=10.5, color=MODEL, bold=True, first=True,
         align=PP_ALIGN.LEFT)
    for t in ["omega        <- WHEEL SPEED          tire_load    <- Fz",
              "torque       <- Trq_T2W              long_force   <- Fx",
              "long_slip    <- kappa                lat_force    <- Fy",
              "lat_slip     <- alpha                tire_friction, normalised forms",
              "",
              "VehicleTelemetryData:  speed  steer  throttle  brake  engine_rpm  gear  drag"]:
        line(tf, t, font=F_MONO, size=10, color=INK, space_before=2, align=PP_ALIGN.LEFT)

    tf = textbox(s, 0.62, 3.74, 6.0, 0.28)
    line(tf, "WHAT IT RETIRES", font=F_DISPLAY, size=13, color=MODEL, bold=True, first=True)
    tf = textbox(s, 0.62, 4.04, 6.1, 1.3)
    for t in ["the vdot inverse -- read torque directly",
              "the unverified 2.52 m^2 drag area -- drag is reported",
              "per-wheel READ needing a fork -- telemetry is per wheel",
              "open items C-1, C-3, C-4, C-5"]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2,
             first=t.startswith("the vdot"), align=PP_ALIGN.LEFT)

    tf = textbox(s, 6.90, 3.74, 6.0, 0.28)
    line(tf, "WHAT IT DOES NOT", font=F_DISPLAY, size=13, color=WARN, bold=True, first=True)
    tf = textbox(s, 6.90, 4.04, 5.9, 1.3)
    for t in ["telemetry is READ-ONLY. Injecting Taxl still goes",
              "through apply_control -- one scalar throttle, split",
              "by CARLA's differential.",
              "So: measure per wheel, command lumped.",
              "Per-wheel COMMAND still needs the SimpleWheeled fork."]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2,
             first=t.startswith("telemetry"), align=PP_ALIGN.LEFT)

    mono_block(s, 0.62, 5.46, 12.1, 0.86, [
        ("Taxl_i -> pedal map -> CARLA -> get_telemetry_data().wheels[i].omega -> w_cmd_i", HW, True),
        ("                                                     .wheels[i].torque -> Trq_T2W_i", HW, True),
    ])

    caption(s, "That is the CarMaker loop, per wheel, with no reconstruction and no fork. "
               "Upgrading is the highest-value action in this study.", y=6.46, color=MODEL)
    return s


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    for fn in (slide_01, slide_02, slide_03, slide_04, slide_05, slide_06):
        fn(prs)
    prs.save(OUT)
    print("wrote %s  (%d slides, %.0f KB)"
          % (OUT.name, len(prs.slides._sldIdLst), OUT.stat().st_size / 1024))


if __name__ == "__main__":
    main()
