#!/usr/bin/env python
"""Regenerate AxleDynoCarMakerCoupling.pptx.

The deck is a build product. This script is its source -- a .pptx does not diff,
so edit here and regenerate rather than editing the binary, unless the change is
one-off presentation polish.

Everything is emitted as native PowerPoint shapes (rectangles, connectors, text
boxes), never as embedded images, so every box and arrow stays editable in
PowerPoint. Content follows AxleDynoCarMakerCoupling.md.

Six slides, in one chain: the loop -> when Trq_T2W[k] is computed within the
integration step -> how the tire model computes it -> what can make the loop
ring -> summary.

    pip install python-pptx
    python doc/xil/make_slides_pptx.py

Colour semantics, shared with the HTML deck and load-bearing throughout:

    HW    warm amber   hardware, measured, physical
    MODEL cool teal    CarMaker, computed
    INK   near-black   our Simulink bypass -- the fulcrum between the two
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

OUT = pathlib.Path(__file__).resolve().parent / "AxleDynoCarMakerCoupling.pptx"

# --------------------------------------------------------------------------
# tokens
# --------------------------------------------------------------------------

INK       = RGBColor(0x11, 0x18, 0x20)
INK_2     = RGBColor(0x4A, 0x57, 0x64)
INK_3     = RGBColor(0x78, 0x83, 0x8F)
RULE      = RGBColor(0xC7, 0xCE, 0xD5)
PAPER     = RGBColor(0xFF, 0xFF, 0xFF)
PANEL     = RGBColor(0xF4, 0xF6, 0xF7)

HW        = RGBColor(0xA9, 0x61, 0x1A)
HW_SOFT   = RGBColor(0xF2, 0xE4, 0xD3)
MODEL     = RGBColor(0x1B, 0x64, 0x6D)
MODEL_SOFT= RGBColor(0xD9, 0xE9, 0xEA)
WARN      = RGBColor(0x93, 0x37, 0x23)
WARN_SOFT = RGBColor(0xF3, 0xDE, 0xD8)

F_DISPLAY = "Bahnschrift"      # Windows' DIN -- the industrial-standard face
F_BODY    = "Segoe UI"
F_MONO    = "Consolas"

DASH = MSO_LINE_DASH_STYLE.DASH

# 16:9 at 13.333 x 7.5 in
SW, SH = Emu(12192000), Emu(6858000)
IN = 914400


def inch(v: float) -> Emu:
    return Emu(int(round(v * IN)))


# --------------------------------------------------------------------------
# primitives
# --------------------------------------------------------------------------

def textbox(slide, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def line(tf, text, *, font=F_BODY, size=12, color=INK, bold=False,
         space_before=0, space_after=0, align=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return p


def box(slide, x, y, w, h, *, fill=None, edge=RULE, width=1.25, dash=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge
    sh.line.width = Pt(width)
    if dash is not None:
        sh.line.dash_style = dash
    sh.text_frame.word_wrap = True
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = inch(0.10)
    tf.margin_top = tf.margin_bottom = inch(0.06)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return sh


def labelled_box(slide, x, y, w, h, head, body, *,
                 fill=None, edge=RULE, head_color=None, width=1.25):
    sh = box(slide, x, y, w, h, fill=fill, edge=edge, width=width)
    tf = sh.text_frame
    line(tf, head, font=F_DISPLAY, size=13, color=head_color or edge,
         bold=True, first=True)
    for i, b in enumerate(body):
        line(tf, b, font=F_MONO, size=9.5, color=INK_2, space_before=3 if i == 0 else 1)
    return sh


def arrow(slide, x1, y1, x2, y2, *, color=INK, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    inch(x1), inch(y1), inch(x2), inch(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash is not None:
        cn.line.dash_style = dash
    tail = cn.line._get_or_add_ln().find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
    if tail is None:
        from pptx.oxml.ns import qn
        tail = cn.line._get_or_add_ln().makeelement(qn("a:tailEnd"), {})
        cn.line._get_or_add_ln().append(tail)
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return cn


def plain_line(slide, x1, y1, x2, y2, *, color=RULE, width=1.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    inch(x1), inch(y1), inch(x2), inch(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash is not None:
        cn.line.dash_style = dash
    return cn


def caption(slide, text, *, y=6.62, color=INK_3, size=9.5, font=F_MONO):
    tf = textbox(slide, 0.62, y, 12.1, 0.22)
    line(tf, text, font=font, size=size, color=color, first=True)


def new_slide(prs, num, eyebrow, title, lede=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER

    plain_line(s, 0.62, 0.52, 12.72, 0.52, color=RULE, width=1.0)

    tf = textbox(s, 0.62, 0.24, 8.0, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "%02d / 06" % num
    r.font.name, r.font.size, r.font.bold = F_MONO, Pt(9), True
    r.font.color.rgb = HW
    r = p.add_run(); r.text = "    " + eyebrow.upper()
    r.font.name, r.font.size = F_MONO, Pt(9)
    r.font.color.rgb = INK_3

    tf = textbox(s, 0.62, 0.70, 12.1, 0.62)
    line(tf, title, font=F_DISPLAY, size=30, color=INK, bold=True, first=True)

    if lede:
        tf = textbox(s, 0.62, 1.36, 11.4, 0.5)
        line(tf, lede, font=F_BODY, size=12.5, color=INK_2, first=True)
    return s


def mono_block(slide, x, y, w, h, rows):
    """A fixed-pitch panel. rows = [(text, color, bold), ...]"""
    sh = box(slide, x, y, w, h, fill=PANEL, edge=RULE, width=1.0)
    tf = sh.text_frame
    tf.margin_left = inch(0.16)
    tf.margin_top = inch(0.10)
    for i, (t, c, b) in enumerate(rows):
        # paragraph 0 of an autoshape defaults to CENTER, so set it explicitly
        line(tf, t, font=F_MONO, size=11, color=c, bold=b,
             space_before=0 if i == 0 else 2, first=(i == 0),
             align=PP_ALIGN.LEFT)
    return sh


# --------------------------------------------------------------------------
# slides
# --------------------------------------------------------------------------

def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = PAPER

    tf = textbox(s, 0.9, 1.15, 8.0, 0.3)
    line(tf, "FIXS  ·  XIL BENCH  ·  dev_v0.9.0", font=F_MONO, size=10,
         color=HW, bold=True, first=True)

    tf = textbox(s, 0.9, 1.62, 11.2, 1.5)
    line(tf, "Axle Dyno ↔ CarMaker Coupling", font=F_DISPLAY, size=52,
         color=INK, bold=True, first=True)

    tf = textbox(s, 0.9, 3.02, 9.6, 1.0)
    line(tf, "How the working implementation exchanges torque and speed with CarMaker "
             "across a single scalar per corner — where the one-step delay comes from, "
             "why the loop is stable, and what the architecture cannot tell you about itself.",
         font=F_BODY, size=14, color=INK_2, first=True)

    mono_block(s, 0.9, 4.16, 8.2, 1.06, [
        ("w[k] = w[k-1] + (K*Ts / I) * ( Taxl[k] + Trq_T2W[k] )", INK, True),
        ("", INK, False),
        ("K = 1     Ts = 0.001 s     I = Wheel.<i>.I  (CarMaker parameter)", INK_2, False),
    ])

    y = 5.56
    for col, lbl in ((HW, "hardware · measured · physical"),
                     (MODEL, "CarMaker · computed"),
                     (INK, "our Simulink bypass")):
        sw = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0.9), inch(y + 0.03),
                                inch(0.15), inch(0.15))
        sw.shadow.inherit = False
        sw.fill.solid(); sw.fill.fore_color.rgb = col
        sw.line.fill.background()
        tf = textbox(s, 1.16, y, 6.0, 0.26)
        line(tf, lbl, font=F_MONO, size=10, color=INK_2, first=True)
        y += 0.30
    return s


def slide_02(prs):
    s = new_slide(prs, 2, "The loop",
                  "What each part owns, and what crosses between them",
                  "Three parts, four signals, one equation. The wheel speed is computed in "
                  "Simulink and sent to both sides. Nothing else crosses.")

    labelled_box(s, 0.62, 2.02, 3.15, 1.70, "HARDWARE",
                 ["real powertrain on the dyno", "",
                  "The dyno speed controller", "makes the shaft follow w."],
                 fill=HW_SOFT, edge=HW, head_color=HW, width=1.5)

    labelled_box(s, 4.72, 1.86, 3.90, 2.02, "SIMULINK — OUR BLOCK",
                 ["w[k] = w[k-1]",
                  "       + (Ts / I) * ( Taxl[k]",
                  "       + Trq_T2W[k] )",
                  "",
                  "Ts = 0.001 s, same step as CarMaker",
                  "I  = Wheel.<i>.I, read from CarMaker"],
                 fill=None, edge=INK, head_color=INK, width=2.5)

    labelled_box(s, 9.57, 2.02, 3.15, 1.70, "CARMAKER",
                 ["tire, body, suspension, road", "",
                  "Uses w to find tire slip,", "then returns the tire torque."],
                 fill=MODEL_SOFT, edge=MODEL, head_color=MODEL, width=1.5)

    # hardware -> simulink
    arrow(s, 3.77, 2.42, 4.66, 2.42, color=HW, width=2.5)
    tf = textbox(s, 3.78, 2.14, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "Taxl[k]", font=F_MONO, size=9.5, color=HW, bold=True, first=True)
    tf = textbox(s, 3.78, 2.48, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "[Nm]", font=F_MONO, size=8.5, color=INK_3, first=True)

    # simulink -> hardware
    arrow(s, 4.66, 3.42, 3.77, 3.42, color=HW, width=2.5, dash=DASH)
    tf = textbox(s, 3.78, 3.14, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "w[k]", font=F_MONO, size=9.5, color=HW, bold=True, first=True)
    tf = textbox(s, 3.78, 3.48, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "[rad/s]", font=F_MONO, size=8.5, color=INK_3, first=True)

    # simulink -> carmaker
    arrow(s, 8.62, 2.42, 9.51, 2.42, color=MODEL, width=2.5)
    tf = textbox(s, 8.63, 2.14, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "w[k]", font=F_MONO, size=9.5, color=MODEL, bold=True, first=True)
    tf = textbox(s, 8.63, 2.48, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "[rad/s]", font=F_MONO, size=8.5, color=INK_3, first=True)

    # carmaker -> simulink
    arrow(s, 9.51, 3.42, 8.62, 3.42, color=MODEL, width=2.5)
    tf = textbox(s, 8.63, 3.14, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "Trq_T2W", font=F_MONO, size=9, color=MODEL, bold=True, first=True)
    tf = textbox(s, 8.63, 3.48, 0.88, 0.22, align=PP_ALIGN.CENTER)
    line(tf, "[Nm]", font=F_MONO, size=8.5, color=INK_3, first=True)

    plain_line(s, 0.62, 4.30, 12.72, 4.30, color=RULE)

    facts = [
        ("The same w[k] goes to both sides.", HW,
         "CarMaker is given the computed speed, not the measured one."),
        ("The real shaft speed stays in the dyno controller.", INK,
         "If the servo lags, CarMaker never finds out."),
        ("CarMaker has no wheel speed of its own.", MODEL,
         "Its integrator went with the PowerTrain module, so ours is the only one."),
        ("CarMaker never uses Taxl either.", INK,
         "We send it, but with the support torques zeroed nothing reads it."),
    ]
    yy = 4.46
    for head, col, tail in facts:
        tf = textbox(s, 0.62, yy, 12.1, 0.24)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = head
        r.font.name, r.font.size, r.font.bold = F_MONO, Pt(10), True
        r.font.color.rgb = col
        r = p.add_run(); r.text = "   " + tail
        r.font.name, r.font.size = F_MONO, Pt(10)
        r.font.color.rgb = INK_2
        yy += 0.30

    caption(s, "Read it as a sentence: add the measured axle torque to the tire torque from CarMaker, "
               "divide by the wheel inertia, integrate.", y=5.82)
    caption(s, "CarMaker does this same integration natively in its DriveLine (RefMan Fig. 16.83). We "
               "rebuilt that block with the same inertia parameter,", y=6.06)
    caption(s, "matching IPG's own UserPowerTrain.mdl. All CarMaker brake torque is zeroed here, because "
               "real braking is already inside Taxl.", y=6.30)
    return s


def slide_03(prs):
    s = new_slide(prs, 3, "When", "One integration step, in a fixed order",
                  "Everything from here on answers one question: how is Trq_T2W[k] produced. "
                  "Start with when CarMaker computes it.")

    # ---- the time axis on the far left: this IS the integration step ----
    TOP, BOT = 2.06, 6.10
    plain_line(s, 1.42, TOP, 12.72, TOP, color=HW, width=1.5)
    plain_line(s, 1.42, BOT, 12.72, BOT, color=HW, width=1.5)
    arrow(s, 1.42, TOP, 1.42, BOT, color=HW, width=2.5)

    tf = textbox(s, 0.62, TOP - 0.13, 0.78, 0.26)
    line(tf, "t = k-1", font=F_MONO, size=10, color=HW, bold=True, first=True)
    tf = textbox(s, 9.30, TOP - 0.28, 3.4, 0.26, align=PP_ALIGN.RIGHT)
    line(tf, "wheel speed = w[k-1]", font=F_MONO, size=10, color=HW, bold=True, first=True)

    tf = textbox(s, 0.62, BOT - 0.13, 0.78, 0.26)
    line(tf, "t = k", font=F_MONO, size=10, color=HW, bold=True, first=True)
    tf = textbox(s, 9.30, BOT - 0.28, 3.4, 0.26, align=PP_ALIGN.RIGHT)
    line(tf, "wheel speed = w[k]", font=F_MONO, size=10, color=HW, bold=True, first=True)

    tf = textbox(s, 0.62, 3.70, 0.72, 1.0)
    line(tf, "ONE", font=F_MONO, size=9, color=HW, bold=True, first=True)
    line(tf, "STEP", font=F_MONO, size=9, color=HW, bold=True)
    line(tf, "1 ms", font=F_MONO, size=9, color=HW)

    # ---- the module sequence ----
    X, WD = 1.90, 5.30
    tf = textbox(s, X, 2.20, 5.0, 0.24)
    line(tf, "Vehicle", font=F_MONO, size=11, color=INK, bold=True, first=True)
    tf = textbox(s, X + 0.30, 2.46, 6.0, 0.24)
    line(tf, "Steering / Suspension K&C / Aerodynamics / Suspension forces",
         font=F_MONO, size=9.5, color=INK_2, first=True)

    tire = box(s, X, 2.80, WD, 0.56, fill=MODEL_SOFT, edge=MODEL, width=2.0)
    tf = tire.text_frame
    line(tf, "Tire", font=F_MONO, size=11, color=MODEL, bold=True, first=True)
    line(tf, "uses w[k-1], emits Trq_T2W[k]", font=F_MONO, size=9.5, color=INK_2, space_before=2)

    tf = textbox(s, X, 3.48, 5.0, 0.24)
    line(tf, "Brake", font=F_MONO, size=11, color=INK, bold=True, first=True)

    pt = box(s, X, 3.80, WD, 0.86, fill=HW_SOFT, edge=HW, width=2.5)
    tf = pt.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "PowerTrain"
    r.font.name, r.font.size, r.font.bold = F_MONO, Pt(11), True
    r.font.color.rgb = INK
    r = p.add_run(); r.text = "     our block"
    r.font.name, r.font.size = F_MONO, Pt(9)
    r.font.color.rgb = HW
    line(tf, "w[k] = w[k-1] + (Ts / I) * ( Taxl[k] + Trq_T2W[k] )",
         font=F_MONO, size=10, color=INK, bold=True, space_before=3)

    bf = box(s, X, 4.82, WD, 0.56, fill=MODEL_SOFT, edge=MODEL, width=2.0)
    tf = bf.text_frame
    line(tf, "Body Frame", font=F_MONO, size=11, color=MODEL, bold=True, first=True)
    line(tf, "chassis, suspension and vehicle speed", font=F_MONO, size=9.5, color=INK_2, space_before=2)

    for y0, y1 in ((3.36, 3.44), (3.72, 3.76), (4.66, 4.78)):
        arrow(s, X + WD / 2, y0, X + WD / 2, y1, color=RULE, width=2.0)

    # ---- the two state advances, called out ----
    arrow(s, X + WD, 4.23, X + WD + 0.42, 4.23, color=HW, width=2.5)
    adv = box(s, X + WD + 0.48, 3.94, 4.58, 0.58, fill=None, edge=HW, width=2.0)
    tf = adv.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    line(tf, "THE WHEEL SPEED ADVANCES HERE", font=F_MONO, size=10.5,
         color=HW, bold=True, first=True)

    arrow(s, X + WD, 5.10, X + WD + 0.42, 5.10, color=MODEL, width=2.0)
    adv2 = box(s, X + WD + 0.48, 4.86, 4.58, 0.48, fill=None, edge=MODEL, width=1.5)
    tf = adv2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    line(tf, "the chassis states advance here", font=F_MONO, size=10,
         color=MODEL, bold=True, first=True)

    caption(s, "The whole sequence is one step of the solver. The wheel speed keeps the value w[k-1] "
               "right up to the PowerTrain box, so the tire, which runs", y=6.34, color=HW)
    caption(s, "earlier, necessarily uses w[k-1]. That is what advancing an ODE state means, not a "
               "delay we added: CarMaker's own powertrain does the same.", y=6.56)

    labelled_box(s, 0.62, 6.80, 5.95, 0.52, "Backward Euler",
                 ["CarMaker's own solver; ours matches. Single stage, one pass."],
                 fill=None, edge=RULE, head_color=INK, width=1.25)
    labelled_box(s, 6.77, 6.80, 5.95, 0.52, "Do not add a second step",
                 ["A Memory block or rate transition in the loop costs damping."],
                 fill=WARN_SOFT, edge=WARN, head_color=WARN, width=1.5)

    return s


def slide_04(prs):
    s = new_slide(prs, 4, "How", "What the tire model does with the wheel speed",
                  "Slip is the gap between our wheel speed and CarMaker's vehicle speed. "
                  "Everything the tire produces follows from that one comparison.")

    labelled_box(s, 0.62, 1.90, 5.55, 0.74, "OUR WHEEL SPEED",
                 ["w[k-1] -> Rim_rotv, from the bench"],
                 fill=HW_SOFT, edge=HW, head_color=HW, width=2.0)
    labelled_box(s, 7.17, 1.90, 5.55, 0.74, "CARMAKER'S VEHICLE SPEED",
                 ["vx at the contact point, plus Fz and friction"],
                 fill=MODEL_SOFT, edge=MODEL, head_color=MODEL, width=1.5)

    plain_line(s, 3.40, 2.64, 3.40, 2.84, color=HW, width=2.5)
    plain_line(s, 9.95, 2.64, 9.95, 2.84, color=MODEL, width=2.0)
    plain_line(s, 3.40, 2.84, 9.95, 2.84, color=RULE, width=2.0)
    arrow(s, 6.67, 2.84, 6.67, 3.02, color=RULE, width=2.5)

    slip = box(s, 1.30, 3.08, 10.74, 0.84, fill=None, edge=RULE, width=1.5)
    tf = slip.text_frame
    p = tf.paragraphs[0]
    for txt, col, bold, sz in (("1   SLIP   ", INK, True, 10.5),
                               ("longitudinal   ", HW, True, 10),
                               ("s = ( w · rBelt_eff  -  vx ) / |…|", INK, False, 10)):
        r = p.add_run(); r.text = txt
        r.font.name, r.font.size, r.font.bold = F_MONO, Pt(sz), bold
        r.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.space_before = Pt(2)
    for txt, col, bold, sz in (("           lateral        ", MODEL, True, 10),
                               ("alpha = atan( vy / |vx| )", INK, False, 10)):
        r = p2.add_run(); r.text = txt
        r.font.name, r.font.size, r.font.bold = F_MONO, Pt(sz), bold
        r.font.color.rgb = col
    line(tf, "how much faster the tread is moving than the road under it, and at what angle",
         font=F_MONO, size=9, color=INK_3, space_before=3)

    arrow(s, 6.67, 3.92, 6.67, 4.08, color=RULE, width=2.5)

    relax = box(s, 1.30, 4.14, 10.74, 0.56, fill=HW_SOFT, edge=HW, width=2.0)
    tf = relax.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    for txt, col, bold, sz in (("2   RELAXATION   ", HW, True, 10.5),
                               ("the carcass takes about 5 cm of rolling to build force, so slip is "
                                "low-passed  ", INK, False, 9.5),
                               ("(the tire's only states)", HW, False, 9)):
        r = p.add_run(); r.text = txt
        r.font.name, r.font.size, r.font.bold = F_MONO, Pt(sz), bold
        r.font.color.rgb = col

    arrow(s, 6.67, 4.70, 6.67, 4.86, color=RULE, width=2.5)

    forces = box(s, 1.30, 4.92, 10.74, 0.84, fill=MODEL_SOFT, edge=MODEL, width=2.0)
    tf = forces.text_frame
    p = tf.paragraphs[0]
    for txt, col, bold, sz in (("3   FORCES   ", MODEL, True, 10.5),
                               ("Fx = Gx(alpha) · Fx0(s)      ", INK, False, 10),
                               ("longitudinal, reduced by cornering", INK_3, False, 9)):
        r = p.add_run(); r.text = txt
        r.font.name, r.font.size, r.font.bold = F_MONO, Pt(sz), bold
        r.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.space_before = Pt(2)
    for txt, col, bold, sz in (("             Fy = Gy(s)     · Fy0(alpha)  ", INK, False, 10),
                               ("lateral, reduced by driving or braking", INK_3, False, 9)):
        r = p2.add_run(); r.text = txt
        r.font.name, r.font.size, r.font.bold = F_MONO, Pt(sz), bold
        r.font.color.rgb = col
    line(tf, "one contact patch, one friction budget, so the two directions take from each other",
         font=F_MONO, size=9, color=INK_3, space_before=3)

    # Fx splits two ways: torque back to us, force into the vehicle body
    arrow(s, 3.60, 5.76, 3.60, 6.04, color=WARN, width=2.5)
    arrow(s, 9.80, 5.76, 9.80, 6.04, color=MODEL, width=2.5)

    t1 = box(s, 1.30, 6.10, 4.60, 0.62, fill=None, edge=WARN, width=2.5)
    tf = t1.text_frame
    line(tf, "4a   TO US", font=F_MONO, size=10, color=WARN, bold=True, first=True)
    line(tf, "Trq_T2W = -r · Fx  ->  our integrator", font=F_MONO, size=10, color=INK, space_before=2)

    t2 = box(s, 7.44, 6.10, 4.60, 0.62, fill=None, edge=MODEL, width=2.5)
    tf = t2.text_frame
    line(tf, "4b   TO THE VEHICLE", font=F_MONO, size=10, color=MODEL, bold=True, first=True)
    line(tf, "Fx pushes the body  ->  CarMaker sets vx", font=F_MONO, size=10, color=INK, space_before=2)

    # the vx return path, drawn round the right-hand side back into stage 1
    plain_line(s, 12.04, 6.41, 12.45, 6.41, color=MODEL, width=2.0, dash=DASH)
    plain_line(s, 12.45, 6.41, 12.45, 3.50, color=MODEL, width=2.0, dash=DASH)
    arrow(s, 12.45, 3.50, 12.08, 3.50, color=MODEL, width=2.0, dash=DASH)
    tf = textbox(s, 12.52, 4.62, 0.78, 0.7)
    line(tf, "vx", font=F_MONO, size=9, color=MODEL, bold=True, first=True)
    line(tf, "next", font=F_MONO, size=9, color=MODEL, bold=True)
    line(tf, "step", font=F_MONO, size=9, color=MODEL, bold=True)

    caption(s, "The bench sets the wheel speed. CarMaker sets the vehicle speed, by integrating these "
               "same tire forces on the vehicle mass at Body Frame.", y=6.90, color=HW)
    caption(s, "Neither one is imposed on the other. They meet only in the slip calculation, which "
               "is the whole coupling.", y=7.12)
    return s


def slide_05(prs):
    s = new_slide(prs, 5, "What can go wrong",
                  "The computed wheel speed can start to oscillate",
                  "This is the one failure mode worth knowing. It shows up in the logs, "
                  "and the dyno is then commanded to chase the wobble.")

    tf = textbox(s, 0.62, 1.94, 6.0, 0.26)
    line(tf, "1  WHAT YOU WOULD SEE", font=F_DISPLAY, size=14, color=WARN, bold=True, first=True)

    tf = textbox(s, 0.62, 2.26, 4.10, 0.80)
    for t in ["rotv and Trq_T2W wobbling at tens of hertz,",
              "on top of the speed the vehicle should be",
              "holding. Because rotv is also the dyno",
              "command, the real shaft is told to chase it."]:
        line(tf, t, font=F_MONO, size=9.5, color=INK_2,
             first=t.startswith("rotv and"))

    tf = textbox(s, 0.62, 3.18, 4.0, 0.24)
    line(tf, "rotv during a slow crawl", font=F_MONO, size=9, color=INK_3, first=True)
    pts = [(0.70, 3.72), (1.50, 3.72), (1.70, 3.66), (1.90, 3.76), (2.10, 3.70), (2.90, 3.72)]
    for a, b in zip(pts, pts[1:]):
        plain_line(s, a[0], a[1], b[0], b[1], color=MODEL, width=2.0)
    tf = textbox(s, 3.00, 3.60, 1.0, 0.24)
    line(tf, "healthy", font=F_MONO, size=9, color=MODEL, first=True)

    rp = [(0.70, 4.34), (1.30, 4.34), (1.42, 4.02), (1.54, 4.64), (1.66, 4.14),
          (1.78, 4.54), (1.90, 4.26), (2.02, 4.44), (2.14, 4.34), (2.90, 4.34)]
    for a, b in zip(rp, rp[1:]):
        plain_line(s, a[0], a[1], b[0], b[1], color=WARN, width=2.0)
    tf = textbox(s, 3.00, 4.22, 1.4, 0.24)
    line(tf, "oscillating", font=F_MONO, size=9, color=WARN, bold=True, first=True)

    plain_line(s, 4.90, 1.90, 4.90, 5.70, color=RULE, width=1.0)

    tf = textbox(s, 5.20, 1.94, 7.5, 0.26)
    line(tf, "2  WHY IT CAN HAPPEN AT ALL", font=F_DISPLAY, size=14, color=INK, bold=True, first=True)
    tf = textbox(s, 5.20, 2.26, 7.5, 0.9)
    line(tf, "The wheel has rotational inertia. The tire carcass has to flex before it can",
         font=F_MONO, size=9.5, color=INK_2, first=True)
    line(tf, "pass force to the road, so it behaves as a spring between the wheel and the",
         font=F_MONO, size=9.5, color=INK_2)
    line(tf, "ground. An inertia on a spring can ring. Every real car has this mode; it is",
         font=F_MONO, size=9.5, color=INK_2)
    line(tf, "the same one behind brake judder and ABS chatter.",
         font=F_MONO, size=9.5, color=INK_2)

    tf = textbox(s, 5.20, 3.46, 7.5, 0.26)
    line(tf, "3  WHY IT ONLY BITES AT LOW SPEED", font=F_DISPLAY, size=14, color=INK, bold=True, first=True)
    tf = textbox(s, 5.20, 3.78, 7.5, 0.9)
    line(tf, "Rolling is what damps it: the tire keeps laying down fresh rubber and renewing",
         font=F_MONO, size=9.5, color=INK_2, first=True)
    line(tf, "its grip. Fast rolling renews it quickly and the ringing dies out. At walking",
         font=F_MONO, size=9.5, color=INK_2)
    line(tf, "pace there is little damping left, and at a standstill none.",
         font=F_MONO, size=9.5, color=INK_2)
    line(tf, "Evaluating the tire once per step removes a further fixed amount, at any speed.",
         font=F_MONO, size=9.5, color=WARN, space_before=4)

    tf = textbox(s, 5.20, 4.98, 7.5, 0.26)
    line(tf, "4  WHY WE DO NOT SEE IT", font=F_DISPLAY, size=14, color=MODEL, bold=True, first=True)
    tf = textbox(s, 5.20, 5.30, 7.5, 0.5)
    line(tf, "Below 2 m/s CarMaker switches to its stand-still tire model, and we put the dyno",
         font=F_MONO, size=9.5, color=INK_2, first=True)
    line(tf, "into torque mode with zero tire torque while stopping. The slow region is never run.",
         font=F_MONO, size=9.5, color=INK_2)

    labelled_box(s, 0.62, 6.06, 3.85, 1.10, "How to check",
                 ["Look at rotv and Trq_T2W through a", "slow crawl. Clean trace, no issue.",
                  "Tens-of-hertz wobble, come back to this."],
                 fill=None, edge=WARN, head_color=WARN, width=1.5)
    labelled_box(s, 4.72, 6.06, 3.85, 1.10, "What would raise the risk",
                 ["A longer sample time, a smaller wheel", "inertia, a shorter relaxation length,",
                  "or a second delay added to the loop."],
                 fill=None, edge=HW, head_color=HW, width=1.25)
    labelled_box(s, 8.82, 6.06, 3.90, 1.10, "The frequency, if you want it",
                 ["Set by wheel inertia against carcass", "stiffness; tens of hertz for a car wheel.",
                  "Our 40-50 Hz is a textbook estimate."],
                 fill=None, edge=RULE, head_color=INK, width=1.25)
    return s


def slide_06(prs):
    s = new_slide(prs, 6, "Summary", "What we can rely on, and what to check",
                  "The loop is sound. Three reasons, then the two things it cannot report on "
                  "itself, then the open items.")

    rely = [("There is only one wheel speed",
             ["Not one in the model and another on the bench.", "Ours is the only one, so the two cannot",
              "disagree."]),
            ("It is IPG's own bypass",
             ["Backward Euler at 1 ms, wheel inertia read", "from CarMaker, block for block the same as",
              "UserPowerTrain.mdl."]),
            ("Errors stay small",
             ["A 2 N·m torque bias settles at about", "0.24 m/s of speed offset. It does not grow:",
              "slip stiffness and drag both push back."])]
    for i, (h, b) in enumerate(rely):
        labelled_box(s, 0.62 + i * 4.10, 1.90, 3.85, 1.30, h, b,
                     fill=MODEL_SOFT, edge=MODEL, head_color=MODEL, width=1.5)

    blind = [("We cannot see the servo error",
              ["CarMaker is given the commanded speed, so it never learns whether the shaft reached",
               "it. If the servo saturates we integrate a torque produced at a different speed, and",
               "nothing warns us.  Log w_act - w_cmd and set an alarm on it."]),
             ("We cannot check Taxl against the model",
              ["CarMaker never reads the axle torque, so a measurement error shows up only as that",
               "small speed offset, which looks like any other modelling choice. Torque calibration",
               "has to be done on the bench."])]
    for i, (h, b) in enumerate(blind):
        labelled_box(s, 0.62 + i * 6.20, 3.42, 5.95, 1.42, h, b,
                     fill=WARN_SOFT, edge=WARN, head_color=WARN, width=1.5)

    rest = [("To check",
             ["Model.USE_MODE in the tire file. 13, 14 or 15", "is expected. It decides both the lateral",
              "coupling and the low-speed behaviour."]),
            ("Recorded simplifications",
             ["All CarMaker brake torque zeroed, since real", "braking is inside Taxl. Support torques",
              "zeroed too: no squat, no mount reaction."]),
            ("Full detail",
             ["doc/xil/AxleDynoCarMakerCoupling.md has the", "step-indexed equations, the source",
              "references and the open-item list."])]
    for i, (h, b) in enumerate(rest):
        labelled_box(s, 0.62 + i * 4.10, 5.06, 3.85, 1.30, h, b,
                     fill=None, edge=RULE, head_color=INK, width=1.25)
    return s


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    for fn in (slide_01, slide_02, slide_03, slide_04, slide_05, slide_06):
        fn(prs)
    prs.save(OUT)
    print("wrote %s  (%d slides, %.0f KB)"
          % (OUT.name, len(prs.slides.__iter__.__self__._sldIdLst), OUT.stat().st_size / 1024))


if __name__ == "__main__":
    main()
